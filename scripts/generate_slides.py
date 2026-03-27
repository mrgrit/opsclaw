#!/usr/bin/env python3
"""
보안 교육과정 lecture.md → PPTX 슬라이드 자동 생성기

각 lecture.md의 ## 헤더를 슬라이드 단위로 분리하고,
내용(텍스트, 코드 블록, 테이블)을 슬라이드에 배치한다.
"""

import os
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ── 스타일 설정 ─────────────────────────────────────────────────────────────
SLIDE_WIDTH = Inches(13.333)  # 16:9
SLIDE_HEIGHT = Inches(7.5)

COLOR_TITLE_BG = RGBColor(0x1B, 0x26, 0x3B)  # 어두운 남색
COLOR_TITLE_TEXT = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_SECTION_BG = RGBColor(0x22, 0x31, 0x4E)
COLOR_BODY_BG = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_BODY_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_CODE_BG = RGBColor(0xF5, 0xF5, 0xF5)
COLOR_CODE_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
COLOR_ACCENT = RGBColor(0x00, 0x96, 0xD6)
FONT_TITLE = "맑은 고딕"
FONT_BODY = "맑은 고딕"
FONT_CODE = "Consolas"


def parse_lecture(filepath: str) -> list[dict]:
    """lecture.md를 파싱하여 슬라이드 데이터 리스트를 반환."""
    with open(filepath, "r", encoding="utf-8") as f:
        content = f.read()

    slides = []
    # # 제목 (Week 01: ...)
    title_match = re.match(r"^#\s+(.+)", content)
    main_title = title_match.group(1) if title_match else Path(filepath).parent.name

    # 학습 목표 추출
    objectives = []
    obj_match = re.search(r"## 학습 목표\n(.*?)(?=\n##|\n---|\Z)", content, re.DOTALL)
    if obj_match:
        for line in obj_match.group(1).strip().split("\n"):
            line = line.strip().lstrip("- ")
            if line:
                objectives.append(line)

    # 타이틀 슬라이드
    slides.append({"type": "title", "title": main_title, "subtitle": ""})

    # 학습 목표 슬라이드
    if objectives:
        slides.append({"type": "objectives", "title": "학습 목표", "items": objectives})

    # ## 헤더 기준으로 분리
    sections = re.split(r"\n(?=## )", content)
    for section in sections:
        lines = section.strip().split("\n")
        if not lines:
            continue

        header_match = re.match(r"^##\s+(.+)", lines[0])
        if not header_match:
            continue

        header = header_match.group(1)
        if header in ("학습 목표", "전제 조건"):
            continue  # 이미 처리했거나 건너뜀

        body_lines = lines[1:]
        body_text = "\n".join(body_lines).strip()

        # 코드 블록 추출
        code_blocks = re.findall(r"```[\w]*\n(.*?)```", body_text, re.DOTALL)

        # 테이블 추출
        tables = []
        table_pattern = re.compile(r"(\|.+\|\n\|[-:| ]+\|\n(?:\|.+\|\n)*)")
        for tm in table_pattern.finditer(body_text):
            tables.append(tm.group(0).strip())

        # 불릿 포인트 추출
        bullets = []
        for line in body_lines:
            line = line.strip()
            if line.startswith("- ") or line.startswith("* "):
                bullets.append(line[2:].strip())
            elif re.match(r"^\d+\.\s", line):
                bullets.append(re.sub(r"^\d+\.\s*", "", line).strip())

        # 하위 섹션 (### 헤더)
        sub_sections = re.findall(r"###\s+(.+)", body_text)

        # 슬라이드 생성 로직
        if "중간고사" in header or "기말" in header:
            slides.append({"type": "section", "title": header})
        elif code_blocks:
            # 코드가 있으면 코드 슬라이드
            summary_bullets = bullets[:5] if bullets else []
            slides.append({
                "type": "content_code",
                "title": header,
                "bullets": summary_bullets,
                "code": code_blocks[0][:600],  # 첫 코드 블록, 600자 제한
            })
        elif tables:
            slides.append({
                "type": "content_table",
                "title": header,
                "table_raw": tables[0],
                "bullets": bullets[:3],
            })
        elif bullets:
            # 불릿이 많으면 분할
            for i in range(0, len(bullets), 6):
                chunk = bullets[i:i+6]
                suffix = f" ({i//6+1})" if len(bullets) > 6 and i > 0 else ""
                slides.append({
                    "type": "content_bullets",
                    "title": header + suffix,
                    "bullets": chunk,
                })
        elif sub_sections:
            slides.append({
                "type": "content_bullets",
                "title": header,
                "bullets": sub_sections[:6],
            })
        else:
            # 텍스트만 있는 경우
            text = body_text[:400].replace("```", "").strip()
            if text:
                slides.append({
                    "type": "content_text",
                    "title": header,
                    "text": text,
                })

    # 마지막: 다음 주 예고
    next_match = re.search(r"## 다음 주 예고\n(.+?)(?:\n##|\Z)", content, re.DOTALL)
    if next_match:
        next_text = next_match.group(1).strip()[:200]
        slides.append({"type": "section", "title": f"다음 주: {next_text.split(chr(10))[0]}"})

    return slides


def add_title_slide(prs: Presentation, data: dict):
    """타이틀 슬라이드."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃

    # 배경
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLOR_TITLE_BG

    # 제목
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = data["title"]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITLE_TEXT
    p.font.name = FONT_TITLE
    p.alignment = PP_ALIGN.CENTER

    # 부제
    if data.get("subtitle"):
        p2 = tf.add_paragraph()
        p2.text = data["subtitle"]
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(0xAA, 0xBB, 0xCC)
        p2.font.name = FONT_BODY
        p2.alignment = PP_ALIGN.CENTER


def add_section_slide(prs: Presentation, data: dict):
    """섹션 구분 슬라이드."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLOR_SECTION_BG

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = data["title"]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITLE_TEXT
    p.font.name = FONT_TITLE
    p.alignment = PP_ALIGN.CENTER


def add_objectives_slide(prs: Presentation, data: dict):
    """학습 목표 슬라이드."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "📌 " + data["title"]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT
    p.font.name = FONT_TITLE

    # 목표 리스트
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, item in enumerate(data["items"][:8]):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = f"✅  {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_BODY_TEXT
        p.font.name = FONT_BODY
        p.space_after = Pt(12)


def add_content_bullets_slide(prs: Presentation, data: dict):
    """불릿 포인트 콘텐츠 슬라이드."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = data["title"]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITLE_BG
    p.font.name = FONT_TITLE

    # 구분선
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.1), Inches(12), Pt(2))  # MSO_SHAPE.RECTANGLE
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT

    # 불릿
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, item in enumerate(data.get("bullets", [])):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        # 마크다운 굵은 글씨 제거
        clean = re.sub(r"\*\*(.+?)\*\*", r"\1", item)
        p.text = f"•  {clean}"
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_BODY_TEXT
        p.font.name = FONT_BODY
        p.space_after = Pt(10)


def add_content_code_slide(prs: Presentation, data: dict):
    """코드 포함 콘텐츠 슬라이드."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = data["title"]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITLE_BG
    p.font.name = FONT_TITLE

    y_offset = 1.2

    # 불릿 (있으면)
    if data.get("bullets"):
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(y_offset), Inches(12), Inches(1.5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for i, item in enumerate(data["bullets"][:3]):
            p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
            clean = re.sub(r"\*\*(.+?)\*\*", r"\1", item)
            p.text = f"•  {clean}"
            p.font.size = Pt(14)
            p.font.color.rgb = COLOR_BODY_TEXT
            p.font.name = FONT_BODY
        y_offset += 1.5

    # 코드 블록
    code_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_offset), Inches(12), Inches(7.5 - y_offset - 0.3))
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x2E)
    tf3 = code_box.text_frame
    tf3.word_wrap = True
    code_text = data.get("code", "")[:500]
    for i, line in enumerate(code_text.split("\n")[:15]):
        p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0xA6, 0xE2, 0x2E)
        p.font.name = FONT_CODE


def add_content_text_slide(prs: Presentation, data: dict):
    """텍스트 콘텐츠 슬라이드."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = data["title"]
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = COLOR_TITLE_BG
    p.font.name = FONT_TITLE

    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", data.get("text", ""))
    p2 = tf2.paragraphs[0]
    p2.text = text
    p2.font.size = Pt(15)
    p2.font.color.rgb = COLOR_BODY_TEXT
    p2.font.name = FONT_BODY


def add_content_table_slide(prs: Presentation, data: dict):
    """테이블 콘텐츠 슬라이드 (불릿으로 대체)."""
    # python-pptx 테이블은 복잡하므로 불릿으로 변환
    table_raw = data.get("table_raw", "")
    rows = [r.strip() for r in table_raw.split("\n") if r.strip() and not r.strip().startswith("|---")]
    items = []
    for row in rows[1:7]:  # 헤더 제외, 6행까지
        cells = [c.strip() for c in row.split("|") if c.strip()]
        if cells:
            items.append(" | ".join(cells[:4]))

    data_conv = {"title": data["title"], "bullets": (data.get("bullets", []) + items)[:8]}
    add_content_bullets_slide(prs, data_conv)


def generate_pptx(lecture_path: str, output_path: str):
    """lecture.md → PPTX 변환."""
    slides_data = parse_lecture(lecture_path)
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    for sd in slides_data:
        stype = sd["type"]
        if stype == "title":
            add_title_slide(prs, sd)
        elif stype == "section":
            add_section_slide(prs, sd)
        elif stype == "objectives":
            add_objectives_slide(prs, sd)
        elif stype == "content_bullets":
            add_content_bullets_slide(prs, sd)
        elif stype == "content_code":
            add_content_code_slide(prs, sd)
        elif stype == "content_text":
            add_content_text_slide(prs, sd)
        elif stype == "content_table":
            add_content_table_slide(prs, sd)

    prs.save(output_path)
    return len(prs.slides)


def main():
    base_dir = Path("/home/opsclaw/opsclaw/security_education")
    output_dir = base_dir / "slides"
    output_dir.mkdir(exist_ok=True)

    courses = sorted([d for d in base_dir.iterdir() if d.is_dir() and d.name.startswith("course")])
    total_slides = 0
    total_files = 0

    for course_dir in courses:
        course_name = course_dir.name
        course_slide_dir = output_dir / course_name
        course_slide_dir.mkdir(exist_ok=True)

        weeks = sorted([d for d in course_dir.iterdir() if d.is_dir() and d.name.startswith("week")])
        for week_dir in weeks:
            lecture_file = week_dir / "lecture.md"
            if not lecture_file.exists():
                continue

            pptx_name = f"{course_name}_{week_dir.name}.pptx"
            pptx_path = course_slide_dir / pptx_name

            try:
                n_slides = generate_pptx(str(lecture_file), str(pptx_path))
                total_slides += n_slides
                total_files += 1
                print(f"  ✅ {pptx_name}: {n_slides} slides")
            except Exception as exc:
                print(f"  ❌ {pptx_name}: {exc}")

    print(f"\n{'='*50}")
    print(f"Total: {total_files} PPTX files, {total_slides} slides")
    print(f"Output: {output_dir}")


if __name__ == "__main__":
    main()
